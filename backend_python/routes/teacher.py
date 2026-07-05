import os
from flask import Blueprint, request, jsonify, current_app
from werkzeug.utils import secure_filename
from database import get_db
from auth import role_required

teacher_bp = Blueprint('teacher', __name__)

def get_teacher_courses(teacher_id):
    conn = get_db()
    c = conn.cursor()
    c.execute("""SELECT co.id,co.name,co.code,co.credits,co.schedule,co.description,
                        COUNT(DISTINCT e.student_id) as enrolled_count
                 FROM courses co
                 LEFT JOIN enrollments e ON e.course_id=co.id
                 WHERE co.teacher_id=? GROUP BY co.id ORDER BY co.name""", (teacher_id,))
    courses = [dict(r) for r in c.fetchall()]
    # Attach sessions to each course
    for co in courses:
        c.execute("SELECT * FROM class_sessions WHERE course_id=? ORDER BY day,start_time", (co['id'],))
        co['sessions'] = [dict(r) for r in c.fetchall()]
    conn.close()
    return courses

# ─── Dashboard ────────────────────────────────────────────────────────────────

@teacher_bp.route('/dashboard', methods=['GET'])
@role_required('Teacher')
def teacher_dashboard():
    conn = get_db()
    c = conn.cursor()
    c.execute("SELECT COUNT(*) FROM courses WHERE teacher_id=?", (request.user_id,))
    my_courses = c.fetchone()[0]
    c.execute("""SELECT COUNT(DISTINCT e.student_id) FROM enrollments e
                 JOIN courses co ON co.id=e.course_id WHERE co.teacher_id=?""", (request.user_id,))
    my_students = c.fetchone()[0]
    c.execute("""SELECT COUNT(*) FROM assignments WHERE teacher_id=?""", (request.user_id,))
    my_assignments = c.fetchone()[0]
    c.execute("""SELECT COUNT(*) FROM submissions sub
                 JOIN assignments a ON a.id=sub.assignment_id WHERE a.teacher_id=?""", (request.user_id,))
    submissions = c.fetchone()[0]
    c.execute("""SELECT COUNT(CASE WHEN att.status='Present' THEN 1 END) as present, COUNT(att.id) as total
                 FROM attendance att JOIN courses co ON co.id=att.course_id WHERE co.teacher_id=?""", (request.user_id,))
    row = c.fetchone()
    att_rate = round((row['present']/(row['total'])*100) if row['total'] else 0, 1)
    c.execute("""SELECT u.name as student_name, a.title as assignment_title,co.code,sub.submitted_at,sub.file_name
                 FROM submissions sub
                 JOIN users u ON u.id=sub.student_id
                 JOIN assignments a ON a.id=sub.assignment_id
                 JOIN courses co ON co.id=a.course_id
                 WHERE a.teacher_id=? ORDER BY sub.submitted_at DESC LIMIT 5""", (request.user_id,))
    recent_submissions = [dict(r) for r in c.fetchall()]
    conn.close()
    return jsonify({'my_courses': my_courses, 'my_students': my_students,
                    'my_assignments': my_assignments, 'submissions': submissions,
                    'attendance_rate': att_rate, 'recent_submissions': recent_submissions}), 200

# ─── Courses ──────────────────────────────────────────────────────────────────

@teacher_bp.route('/courses', methods=['GET'])
@role_required('Teacher')
def teacher_courses():
    return jsonify(get_teacher_courses(request.user_id)), 200

# ─── Class Sessions ───────────────────────────────────────────────────────────

@teacher_bp.route('/courses/<int:cid>/sessions', methods=['GET'])
@role_required('Teacher')
def get_course_sessions(cid):
    conn = get_db()
    c = conn.cursor()
    # Verify teacher owns this course
    c.execute("SELECT id FROM courses WHERE id=? AND teacher_id=?", (cid, request.user_id))
    if not c.fetchone():
        conn.close()
        return jsonify({'message': 'Not found'}), 404
    c.execute("SELECT * FROM class_sessions WHERE course_id=? ORDER BY day,start_time", (cid,))
    sessions = [dict(r) for r in c.fetchall()]
    conn.close()
    return jsonify(sessions), 200

# ─── Students in my courses ────────────────────────────────────────────────────

@teacher_bp.route('/students', methods=['GET'])
@role_required('Teacher')
def teacher_students():
    course_id = request.args.get('course_id')
    conn = get_db()
    c = conn.cursor()
    if course_id:
        c.execute("""SELECT u.id,u.name,u.email,u.student_id,u.major,u.phone,
                            co.name as course_name, co.code
                     FROM users u
                     JOIN enrollments e ON e.student_id=u.id
                     JOIN courses co ON co.id=e.course_id
                     WHERE co.teacher_id=? AND co.id=? ORDER BY u.name""", (request.user_id, course_id))
    else:
        c.execute("""SELECT DISTINCT u.id,u.name,u.email,u.student_id,u.major,u.phone
                     FROM users u
                     JOIN enrollments e ON e.student_id=u.id
                     JOIN courses co ON co.id=e.course_id
                     WHERE co.teacher_id=? ORDER BY u.name""", (request.user_id,))
    students = [dict(r) for r in c.fetchall()]
    conn.close()
    return jsonify(students), 200

# ─── Attendance ────────────────────────────────────────────────────────────────

@teacher_bp.route('/attendance', methods=['GET'])
@role_required('Teacher')
def get_attendance():
    course_id = request.args.get('course_id')
    date = request.args.get('date')
    conn = get_db()
    c = conn.cursor()
    if course_id and date:
        c.execute("""SELECT u.id as student_id, u.name, u.student_id as sid,
                            COALESCE(a.status,'') as status
                     FROM users u
                     JOIN enrollments e ON e.student_id=u.id AND e.course_id=?
                     LEFT JOIN attendance a ON a.student_id=u.id AND a.course_id=? AND a.date=?
                     ORDER BY u.name""", (course_id, course_id, date))
    else:
        c.execute("""SELECT a.*,u.name as student_name,co.name as course_name
                     FROM attendance a
                     JOIN users u ON u.id=a.student_id
                     JOIN courses co ON co.id=a.course_id
                     WHERE co.teacher_id=? ORDER BY a.date DESC""", (request.user_id,))
    records = [dict(r) for r in c.fetchall()]
    conn.close()
    return jsonify(records), 200

@teacher_bp.route('/attendance', methods=['POST'])
@role_required('Teacher')
def save_attendance():
    data = request.get_json()
    course_id = data.get('course_id')
    date = data.get('date')
    records = data.get('records', [])
    conn = get_db()
    for r in records:
        conn.execute("""INSERT INTO attendance (student_id,course_id,date,status)
                        VALUES (?,?,?,?)
                        ON CONFLICT(student_id,course_id,date) DO UPDATE SET status=excluded.status""",
                     (r['student_id'], course_id, date, r['status']))
    conn.commit()
    conn.close()
    return jsonify({'message': 'Attendance saved'}), 200

# ─── Grades ────────────────────────────────────────────────────────────────────

@teacher_bp.route('/grades', methods=['GET'])
@role_required('Teacher')
def get_grades():
    course_id = request.args.get('course_id')
    semester = request.args.get('semester', 'Spring 2026')
    conn = get_db()
    c = conn.cursor()
    if course_id:
        c.execute("""SELECT u.id as student_id, u.name, u.student_id as sid,
                            COALESCE(g.exam_score,'') as exam_score,
                            COALESCE(g.score,'') as score,
                            COALESCE(g.grade_letter,'') as grade_letter,
                            COALESCE(g.semester,?) as semester
                     FROM users u JOIN enrollments e ON e.student_id=u.id AND e.course_id=?
                     LEFT JOIN grades g ON g.student_id=u.id AND g.course_id=? AND g.semester=?
                     ORDER BY u.name""", (semester, course_id, course_id, semester))
        records = [dict(r) for r in c.fetchall()]

        # Enrich with live attendance % and assignment avg per student
        for r in records:
            sid = r['student_id']
            c.execute("""SELECT COUNT(CASE WHEN status='Present' THEN 1 END) as present,
                                COUNT(*) as total
                         FROM attendance WHERE student_id=? AND course_id=?""", (sid, course_id))
            att = c.fetchone()
            r['attendance_pct'] = round((att['present'] / att['total'] * 100) if att['total'] else 0, 1)

            c.execute("""SELECT ROUND(AVG(sub.score),1) as avg_score
                         FROM submissions sub
                         JOIN assignments a ON a.id=sub.assignment_id
                         WHERE sub.student_id=? AND a.course_id=? AND sub.score IS NOT NULL""",
                      (sid, course_id))
            asgn = c.fetchone()
            r['assignment_avg'] = asgn['avg_score']

            # Compute composite
            exam_s = r.get('exam_score')
            try:
                exam_f = float(exam_s) if exam_s not in ('', None) else None
            except (ValueError, TypeError):
                exam_f = None
            asgn_f = float(r['assignment_avg']) if r['assignment_avg'] is not None else 0.0
            att_f = float(r['attendance_pct'])

            if exam_f is not None:
                r['composite'] = round(att_f * 0.10 + asgn_f * 0.40 + exam_f * 0.50, 1)
            else:
                r['composite'] = None
    else:
        c.execute("""SELECT g.*,u.name as student_name,co.name as course_name
                     FROM grades g JOIN users u ON u.id=g.student_id
                     JOIN courses co ON co.id=g.course_id
                     WHERE co.teacher_id=? ORDER BY co.name,u.name""", (request.user_id,))
        records = [dict(r) for r in c.fetchall()]
    conn.close()
    return jsonify(records), 200

@teacher_bp.route('/grades', methods=['POST'])
@role_required('Teacher')
def save_grades():
    data = request.get_json()
    course_id = data.get('course_id')
    semester = data.get('semester', 'Spring 2026')
    records = data.get('records', [])

    def letter(s):
        if s is None: return ''
        s = float(s)
        if s >= 95: return 'A+'
        if s >= 90: return 'A'
        if s >= 85: return 'B+'
        if s >= 80: return 'B'
        if s >= 75: return 'C+'
        if s >= 70: return 'C'
        if s >= 60: return 'D'
        return 'F'

    conn = get_db()
    c = conn.cursor()
    for r in records:
        exam_score = r.get('exam_score')
        exam_f = float(exam_score) if exam_score not in ('', None) else None

        # Get live attendance pct
        c.execute("""SELECT COUNT(CASE WHEN status='Present' THEN 1 END) as present, COUNT(*) as total
                     FROM attendance WHERE student_id=? AND course_id=?""",
                  (r['student_id'], course_id))
        att = c.fetchone()
        att_pct = round((att['present'] / att['total'] * 100) if att['total'] else 0, 1)

        # Get live assignment avg
        c.execute("""SELECT ROUND(AVG(sub.score),1) as avg_score
                     FROM submissions sub
                     JOIN assignments a ON a.id=sub.assignment_id
                     WHERE sub.student_id=? AND a.course_id=? AND sub.score IS NOT NULL""",
                  (r['student_id'], course_id))
        asgn = c.fetchone()
        asgn_avg = float(asgn['avg_score']) if asgn['avg_score'] is not None else 0.0

        if exam_f is not None:
            composite = round(att_pct * 0.10 + asgn_avg * 0.40 + exam_f * 0.50, 1)
        else:
            composite = None

        gl = letter(composite) if composite is not None else ''

        conn.execute("""INSERT INTO grades (student_id,course_id,semester,score,grade_letter,exam_score,attendance_score,assignment_score)
                        VALUES (?,?,?,?,?,?,?,?)
                        ON CONFLICT(student_id,course_id,semester) DO UPDATE SET
                            score=excluded.score, grade_letter=excluded.grade_letter,
                            exam_score=excluded.exam_score,
                            attendance_score=excluded.attendance_score,
                            assignment_score=excluded.assignment_score""",
                     (r['student_id'], course_id, semester, composite, gl,
                      exam_f, att_pct, asgn_avg if asgn_avg else None))

        # Notify student
        if composite is not None:
            conn.execute("INSERT INTO notifications (user_id, message) VALUES (?,?)",
                         (r['student_id'], f"Your grade has been updated: {composite}/100 ({gl})"))
    conn.commit()
    conn.close()
    return jsonify({'message': 'Grades saved'}), 200

# ─── Assignments ───────────────────────────────────────────────────────────────

@teacher_bp.route('/assignments', methods=['GET'])
@role_required('Teacher')
def get_assignments():
    conn = get_db()
    c = conn.cursor()
    c.execute("""SELECT a.*,co.name as course_name,co.code,
                        COUNT(DISTINCT sub.id) as submission_count
                 FROM assignments a JOIN courses co ON co.id=a.course_id
                 LEFT JOIN submissions sub ON sub.assignment_id=a.id
                 WHERE a.teacher_id=? GROUP BY a.id ORDER BY a.deadline ASC""", (request.user_id,))
    assignments = [dict(r) for r in c.fetchall()]

    for assn in assignments:
        c.execute("""SELECT sub.*,u.name as student_name,u.student_id as sid
                     FROM submissions sub JOIN users u ON u.id=sub.student_id
                     WHERE sub.assignment_id=?""", (assn['id'],))
        assn['submissions'] = [dict(r) for r in c.fetchall()]
    conn.close()
    return jsonify(assignments), 200

@teacher_bp.route('/assignments', methods=['POST'])
@role_required('Teacher')
def add_assignment():
    # Handle both JSON and Form Data
    if request.is_json:
        data = request.get_json()
        file = None
    else:
        data = request.form
        file = request.files.get('file')

    title = data.get('title')
    course_id = data.get('course_id')

    if not title or not course_id:
        return jsonify({'message': 'Title and Course ID are required'}), 400

    file_path = None
    file_name = None

    if file and file.filename != '':
        filename = secure_filename(file.filename)
        # Unique name to avoid collisions
        unique_name = f"assn_{int(os.getpid())}_{filename}"
        file.save(os.path.join(current_app.config['UPLOAD_FOLDER'], unique_name))
        file_path = f"/api/uploads/{unique_name}"
        file_name = filename

    conn = get_db()
    c = conn.cursor()
    c.execute("""INSERT INTO assignments (course_id, teacher_id, title, description, deadline, file_path, file_name)
                 VALUES (?,?,?,?,?,?,?)""",
              (course_id, request.user_id, title, data.get('description'), data.get('deadline'), file_path, file_name))
    conn.commit()
    assn_id = c.lastrowid
    
    c.execute("SELECT student_id FROM enrollments WHERE course_id=?", (course_id,))
    students = c.fetchall()
    for row in students:
        conn.execute("INSERT INTO notifications (user_id, message) VALUES (?,?)",
                     (row['student_id'], f"New assignment posted: {title}"))
    conn.commit()
    conn.close()
    return jsonify({'message': 'Assignment created', 'id': assn_id}), 201

@teacher_bp.route('/assignments/<int:aid>', methods=['PUT'])
@role_required('Teacher')
def update_assignment(aid):
    data = request.get_json()
    conn = get_db()
    conn.execute("""UPDATE assignments SET title=?,description=?,deadline=? WHERE id=? AND teacher_id=?""",
                 (data.get('title'), data.get('description'), data.get('deadline'), aid, request.user_id))
    conn.commit()
    conn.close()
    return jsonify({'message': 'Assignment updated'}), 200

@teacher_bp.route('/assignments/<int:aid>', methods=['DELETE'])
@role_required('Teacher')
def delete_assignment(aid):
    conn = get_db()
    conn.execute("DELETE FROM assignments WHERE id=? AND teacher_id=?", (aid, request.user_id))
    conn.commit()
    conn.close()
    return jsonify({'message': 'Assignment deleted'}), 200

@teacher_bp.route('/assignments/<int:aid>/grade', methods=['POST'])
@role_required('Teacher')
def grade_submission(aid):
    data = request.get_json()
    student_id = data.get('student_id')
    score = data.get('score')
    conn = get_db()
    conn.execute("UPDATE submissions SET score=? WHERE assignment_id=? AND student_id=?", (score, aid, student_id))
    conn.commit()
    conn.execute("INSERT INTO notifications (user_id,message) VALUES (?,?)",
                 (student_id, f"Your assignment has been graded: {score}/100"))
    conn.commit()
    conn.close()
    return jsonify({'message': 'Graded'}), 200

# ─── Announcements ────────────────────────────────────────────────────────────

@teacher_bp.route('/announcements', methods=['GET'])
@role_required('Teacher')
def get_announcements():
    conn = get_db()
    c = conn.cursor()
    c.execute("""SELECT a.*,co.name as course_name,co.code
                 FROM announcements a LEFT JOIN courses co ON co.id=a.course_id
                 WHERE a.teacher_id=? ORDER BY a.created_at DESC""", (request.user_id,))
    announcements = [dict(r) for r in c.fetchall()]
    conn.close()
    return jsonify(announcements), 200

@teacher_bp.route('/announcements', methods=['POST'])
@role_required('Teacher')
def add_announcement():
    data = request.get_json()
    for f in ['title','content']:
        if not data.get(f):
            return jsonify({'message': f'{f} is required'}), 400
    conn = get_db()
    c = conn.cursor()
    c.execute("INSERT INTO announcements (teacher_id,course_id,title,content) VALUES (?,?,?,?)",
              (request.user_id, data.get('course_id'), data['title'], data['content']))
    conn.commit()
    ann_id = c.lastrowid
    if data.get('course_id'):
        c.execute("SELECT student_id FROM enrollments WHERE course_id=?", (data['course_id'],))
        for row in c.fetchall():
            conn.execute("INSERT INTO notifications (user_id,message) VALUES (?,?)",
                         (row['student_id'], f"New announcement: {data['title']}"))
        conn.commit()
    conn.close()
    return jsonify({'message': 'Announcement posted', 'id': ann_id}), 201

@teacher_bp.route('/announcements/<int:aid>', methods=['PUT'])
@role_required('Teacher')
def update_announcement(aid):
    data = request.get_json()
    conn = get_db()
    conn.execute("UPDATE announcements SET title=?,content=?,course_id=? WHERE id=? AND teacher_id=?",
                 (data.get('title'), data.get('content'), data.get('course_id'), aid, request.user_id))
    conn.commit()
    conn.close()
    return jsonify({'message': 'Updated'}), 200

@teacher_bp.route('/announcements/<int:aid>', methods=['DELETE'])
@role_required('Teacher')
def delete_announcement(aid):
    conn = get_db()
    conn.execute("DELETE FROM announcements WHERE id=? AND teacher_id=?", (aid, request.user_id))
    conn.commit()
    conn.close()
    return jsonify({'message': 'Deleted'}), 200

# ─── AI Taxonomy Analysis ──────────────────────────────────────────────────────

@teacher_bp.route('/taxonomy/analyze', methods=['POST'])
@role_required('Teacher')
def analyze_taxonomy():
    """
    Accepts a .pptx file, parses slide text with python-pptx, then calls
    OpenAI to classify each slide by Bloom's Taxonomy level, generate
    improvement suggestions for the teacher, and generate student questions.
    """
    if 'file' not in request.files:
        return jsonify({'message': 'No file provided'}), 400

    file = request.files['file']
    if not file.filename or not file.filename.lower().endswith('.pptx'):
        return jsonify({'message': 'Only .pptx files are supported'}), 400

    # ── Parse PPTX ──
    try:
        import io
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file.read()))
    except Exception as e:
        return jsonify({'message': f'Could not parse PPTX: {str(e)}'}), 400

    slides_data = []
    for i, slide in enumerate(prs.slides):
        title = ''
        bullets = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for j, para in enumerate(shape.text_frame.paragraphs):
                text = para.text.strip()
                if not text or len(text) < 2:
                    continue
                if not title and len(text) > 2:
                    title = text
                elif text != title:
                    bullets.append(text)

        if not title:
            title = f'Slide {i + 1}'

        slides_data.append({
            'id': i + 1,
            'number': i + 1,
            'title': title,
            'bullets': list(dict.fromkeys(bullets))[:8]  # deduplicate, max 8
        })

    if not slides_data:
        return jsonify({'message': 'No readable slides found in this file'}), 400

    # ── Build prompt ──
    slides_text = ''
    for s in slides_data:
        slides_text += f'\nSlide {s["number"]}: {s["title"]}\n'
        for b in s['bullets']:
            slides_text += f'  - {b}\n'

    system_prompt = """You are an expert instructional designer specializing in Bloom's Taxonomy.
Your task is to analyze presentation slides and return a JSON response with the exact structure requested.
Be concise, practical, and actionable in your suggestions."""

    user_prompt = f"""Analyze these presentation slides and for each slide provide:
1. Bloom's Taxonomy classification (one of: remember, understand, apply, analyze, evaluate, create)
2. AI classification reasoning (1-2 sentences)
3. 2-3 specific, actionable improvement suggestions to better align with that level
4. A short, engaging explanation of the slide directly addressing the student, explicitly utilizing the assigned taxonomy level's framing (e.g., if 'Apply', explain how they would apply this; if 'Analyze', explain how to break it down).
5. 2-3 student questions at that taxonomy level to help students understand the content

Bloom's levels reference:
- remember: Recall, List, Name, Define, Recognize
- understand: Explain, Summarize, Classify, Paraphrase
- apply: Demonstrate, Solve, Execute, Implement
- analyze: Differentiate, Compare, Examine, Deconstruct
- evaluate: Judge, Critique, Justify, Defend
- create: Design, Build, Construct, Compose

SLIDES:
{slides_text}

Respond ONLY with valid JSON in this exact structure:
{{
  "slides": [
    {{
      "id": 1,
      "taxonomyLevel": "remember",
      "aiNotes": "Brief reason why this taxonomy level was assigned.",
      "suggestions": [
        "Actionable improvement suggestion 1",
        "Actionable improvement suggestion 2"
      ],
      "studentExplanation": "Engaging explanation of the slide for the student using the taxonomy framing.",
      "studentQuestions": [
        {{"question": "Question that helps students understand?", "answer": "Concise answer/hint for the student."}},
        {{"question": "Another learning question?", "answer": "Answer hint."}}
      ]
    }}
  ]
}}"""

    # ── Call OpenAI ──
    try:
        import json
        from openai import OpenAI
        client = OpenAI(api_key=os.getenv('OPENAI_API_KEY'))

        response = client.chat.completions.create(
            model='gpt-4o-mini',
            messages=[
                {'role': 'system', 'content': system_prompt},
                {'role': 'user',   'content': user_prompt}
            ],
            response_format={'type': 'json_object'},
            temperature=0.3,
            max_tokens=4000
        )

        ai_result = json.loads(response.choices[0].message.content)
    except Exception as e:
        return jsonify({'message': f'AI analysis failed: {str(e)}'}), 500

    # ── Merge AI result with parsed slide data ──
    ai_map = {s['id']: s for s in ai_result.get('slides', [])}
    valid_levels = {'remember', 'understand', 'apply', 'analyze', 'evaluate', 'create'}

    result_slides = []
    for s in slides_data:
        ai = ai_map.get(s['id'], {})
        level = ai.get('taxonomyLevel', 'remember')
        if level not in valid_levels:
            level = 'remember'
        result_slides.append({
            **s,
            'taxonomyLevel': level,
            'aiNotes': ai.get('aiNotes', ''),
            'suggestions': ai.get('suggestions', []),
            'studentExplanation': ai.get('studentExplanation', ''),
            'studentQuestions': ai.get('studentQuestions', [])
        })

    return jsonify({'slides': result_slides}), 200

@teacher_bp.route('/taxonomy/analyze-text', methods=['POST'])
@role_required('Teacher')
def analyze_taxonomy_text():
    """
    Accepts JSON with parsed slide data, calls OpenAI to classify each slide 
    by Bloom's Taxonomy level, generate improvement suggestions for the teacher, 
    and generate student questions.
    """
    data = request.json
    slides_data = data.get('slides', [])
    
    if not slides_data:
        return jsonify({'message': 'No slides provided'}), 400

    # ── Build prompt ──
    slides_text = ''
    for s in slides_data:
        slides_text += f'\nSlide {s.get("number", "?")}: {s.get("title", "")}\n'
        for b in s.get('bullets', []):
            slides_text += f'  - {b}\n'

    system_prompt = """You are an expert instructional designer specializing in Bloom's Taxonomy.
Your task is to analyze presentation slides and return a JSON response with the exact structure requested.
Be concise, practical, and actionable in your suggestions."""

    user_prompt = f"""Analyze these presentation slides and for each slide provide:
1. Bloom's Taxonomy classification (one of: remember, understand, apply, analyze, evaluate, create)
2. AI classification reasoning (1-2 sentences)
3. 2-3 specific, actionable improvement suggestions to better align with that level
4. A short, engaging explanation of the slide directly addressing the student, explicitly utilizing the assigned taxonomy level's framing (e.g., if 'Apply', explain how they would apply this; if 'Analyze', explain how to break it down).
5. 2-3 student questions at that taxonomy level to help students understand the content

Bloom's levels reference:
- remember: Recall, List, Name, Define, Recognize
- understand: Explain, Summarize, Classify, Paraphrase
- apply: Demonstrate, Solve, Execute, Implement
- analyze: Differentiate, Compare, Examine, Deconstruct
- evaluate: Judge, Critique, Justify, Defend
- create: Design, Build, Construct, Compose

SLIDES:
{slides_text}

Respond ONLY with valid JSON in this exact structure:
{{
  "slides": [
    {{
      "id": 1,
      "taxonomyLevel": "remember",
      "aiNotes": "Brief reason why this taxonomy level was assigned.",
      "suggestions": [
        "Actionable improvement suggestion 1",
        "Actionable improvement suggestion 2"
      ],
      "studentExplanation": "Engaging explanation of the slide for the student using the taxonomy framing.",
      "studentQuestions": [
        {{"question": "Question that helps students understand?", "answer": "Concise answer/hint for the student."}},
        {{"question": "Another learning question?", "answer": "Answer hint."}}
      ]
    }}
  ]
}}"""

    # ── Call OpenAI ──
    try:
        import json
        from openai import OpenAI
        client = OpenAI(api_key=os.getenv('OPENAI_API_KEY'))

        response = client.chat.completions.create(
            model='gpt-4o-mini',
            messages=[
                {'role': 'system', 'content': system_prompt},
                {'role': 'user',   'content': user_prompt}
            ],
            response_format={'type': 'json_object'},
            temperature=0.3,
            max_tokens=4000
        )

        ai_result = json.loads(response.choices[0].message.content)
    except Exception as e:
        return jsonify({'message': f'AI analysis failed: {str(e)}'}), 500

    # ── Map result back to slides ──
    ai_map = {s['id']: s for s in ai_result.get('slides', [])}
    valid_levels = {'remember', 'understand', 'apply', 'analyze', 'evaluate', 'create'}

    result_slides = []
    for s in slides_data:
        ai = ai_map.get(s['id'], {})
        level = ai.get('taxonomyLevel', 'remember')
        if level not in valid_levels:
            level = 'remember'
        result_slides.append({
            **s,
            'taxonomyLevel': level,
            'aiNotes': ai.get('aiNotes', ''),
            'suggestions': ai.get('suggestions', []),
            'studentExplanation': ai.get('studentExplanation', ''),
            'studentQuestions': ai.get('studentQuestions', [])
        })

    return jsonify({'slides': result_slides}), 200

import uuid
from pypdf import PdfReader

@teacher_bp.route('/taxonomy/upload-pdf', methods=['POST'])
@role_required('Teacher')
def upload_pdf():
    """
    Accepts a .pdf file, saves it to the UPLOAD_FOLDER with a unique name,
    extracts raw text per page using pypdf, and returns the pdfUrl and parsed slide text.
    """
    if 'file' not in request.files:
        return jsonify({'message': 'No file provided'}), 400

    file = request.files['file']
    if not file.filename or not file.filename.lower().endswith('.pdf'):
        return jsonify({'message': 'Only .pdf files are supported'}), 400

    # Save the file
    from flask import current_app
    import os
    
    unique_filename = f"{uuid.uuid4().hex}_{file.filename}"
    upload_path = os.path.join(current_app.config['UPLOAD_FOLDER'], unique_filename)
    file.save(upload_path)
    
    pdf_url = f"/api/uploads/{unique_filename}"

    # Extract text per page
    try:
        reader = PdfReader(upload_path)
        slides_data = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            # For PDF we just pass the raw text as a single "bullet" or "title" to the AI, 
            # the AI is smart enough to structure it.
            # We split by newlines just to simulate bullets for the prompt builder if needed.
            lines = [line.strip() for line in text.split('\n') if len(line.strip()) > 2]
            
            title = lines[0] if lines else f"Page {i+1}"
            bullets = list(dict.fromkeys(lines[1:]))[:8] if len(lines) > 1 else []
            
            slides_data.append({
                'id': i + 1,
                'number': i + 1,
                'title': title,
                'bullets': bullets,
                'taxonomyLevel': None,
                'suggestions': [],
                'studentQuestions': [],
                'studentExplanation': '',
                'aiNotes': ''
            })
    except Exception as e:
        # If extraction fails, we still return the url so they can see the PDF at least
        return jsonify({'message': f'Failed to extract text: {str(e)}', 'pdfUrl': pdf_url, 'slides': []}), 200

    return jsonify({'pdfUrl': pdf_url, 'slides': slides_data}), 200

